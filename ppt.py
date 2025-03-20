from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Function to add a slide with title and content
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.shapes.placeholders[1]
    
    title_placeholder.text = title
    content_placeholder.text = content

# Function to add a slide with title and bullet points
def add_bulleted_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Title and Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.shapes.placeholders[1]
    
    title_placeholder.text = title
    text_frame = content_placeholder.text_frame
    text_frame.clear()  # Clear default text
    
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point

# Add Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AI Doctor: Disease Prediction Using Machine Learning"
subtitle.text = "Predicting Diseases Based on Symptoms\nYour Name\nDate"

# Add Introduction Slide
add_slide("Introduction to AI Doctor", "AI Doctor is a machine learning project to predict diseases based on symptoms provided by users.\n\nThe system helps healthcare professionals in diagnosing diseases more efficiently.")

# Add Dataset Description Slide
add_slide("Dataset and Data Preprocessing", "The dataset includes diseases, their associated symptoms, and disease occurrence counts.\n\nThe dataset is cleaned by handling missing values, transforming symptoms into numerical format, and encoding them for machine learning.")

# Add Feature Engineering Slide
add_slide("Feature Engineering for Disease Prediction", "We process raw symptom data into numerical features using One-Hot Encoding.\n\nOne-Hot Encoding converts categorical symptom data into usable numerical values for the machine learning model.")

# Add Machine Learning Model Slide
add_slide("Machine Learning Model - Decision Tree Classifier", "We use the Decision Tree Classifier to predict diseases.\n\nThe model learns decision rules from the symptom data to predict the disease.\n\nThe decision tree is trained on the processed dataset.")

# Add Training and Evaluation Slide
add_slide("Model Training and Evaluation", "The model is trained on the cleaned data using a Decision Tree Classifier.\n\nWe evaluate the model using metrics like accuracy, precision, recall, and F1 score.")

# Add Results and Visualizations Slide
add_slide("Results and Decision Tree Visualization", "We evaluate the model’s performance through quantitative metrics and visualizations.\n\nA decision tree visualization helps understand how the model makes predictions.")

# Add Model Deployment Slide
add_slide("Model Deployment and User Interface", "The trained model is saved and deployed for real-world use.\n\nUsers can input symptoms and receive a disease prediction via a simple user interface.")

# Add Potential Applications Slide
add_bulleted_slide("Potential Applications of AI Doctor", [
    "Healthcare Assistance: Assists healthcare professionals in diagnosis.",
    "Symptom Checker: Allows users to check their symptoms and get possible disease predictions.",
    "Early Detection: Helps in identifying early signs of diseases for timely intervention."
])

# Add Conclusion Slide
add_slide("Conclusion and Summary", "AI Doctor is a powerful tool for predicting diseases based on symptoms.\n\nThe system helps in diagnosis, reduces time for doctors, and increases accessibility to healthcare.\n\nFuture scope: integrating more diseases, improving model accuracy, and deploying in real-time systems.")

# Add Project Contributions Slide
add_slide("Key Contributions", "Key contributions include data collection, model training, feature engineering, and evaluation.\n\nThe project aims to improve healthcare diagnosis using AI and machine learning techniques.")

# Add References Slide
add_slide("References", "List all references used in the project, including datasets, research papers, libraries, and tools.")

# Save the presentation
prs.save("AI_Doctor_Presentation.pptx")

print("Presentation generated successfully!")
